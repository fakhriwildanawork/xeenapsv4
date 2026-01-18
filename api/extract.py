
import io
import re
import requests
from readability import Document as ReadabilityDoc
from bs4 import BeautifulSoup
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
import openpyxl
from flask import Flask, request, jsonify

app = Flask(__name__)

# Mengatur batas maksimal konten agar Flask menerima file hingga 25MB
app.config['MAX_CONTENT_LENGTH'] = 25 * 1024 * 1024

def clean_text(text):
    # Pastikan input adalah string untuk menghindari error "bytes-like object"
    if isinstance(text, bytes):
        text = text.decode('utf-8', errors='ignore')
    if not isinstance(text, str):
        text = str(text)
        
    # Membersihkan karakter aneh dan spasi berlebih
    text = re.sub(r'([A-Z])\s(?=[a-z])', r'\1', text) 
    return " ".join(text.split())

def extract_doi(url):
    """Mendeteksi DOI dari URL menggunakan Regex."""
    doi_pattern = r'(10\.\d{4,9}/[-._;()/:A-Z0-9]+)'
    match = re.search(doi_pattern, url, re.IGNORECASE)
    return match.group(1) if match else None

def fetch_from_crossref(doi):
    """Mengambil metadata dari Crossref API (Bebas blokir untuk jurnal)."""
    try:
        api_url = f"https://api.crossref.org/works/{doi}"
        headers = {'User-Agent': 'XeenapsPKM/1.0 (mailto:admin@xeenaps.com)'}
        response = requests.get(api_url, headers=headers, timeout=10)
        if response.status_code == 200:
            data = response.json().get('message', {})
            
            # Format Authors
            authors = [f"{a.get('given', '')} {a.get('family', '')}".strip() for a in data.get('author', [])]
            
            # Extract Year
            year = ""
            pub_date = data.get('published-print') or data.get('published-online') or data.get('issued')
            if pub_date and pub_date.get('date-parts'):
                year = str(pub_date['date-parts'][0][0])

            return {
                "title": data.get('title', [""])[0],
                "authors": authors,
                "year": year,
                "publisher": data.get('container-title', [""])[0] or data.get('publisher', ""),
                "category": "Original Research",
                "type": "Literature"
            }
    except Exception as e:
        print(f"Crossref Error: {e}")
    return None

def extract_metadata_heuristics(full_text, filename_or_title):
    # Pastikan input adalah string
    text_str = str(full_text)
        
    metadata = {
        "title": filename_or_title.rsplit('.', 1)[0].replace("_", " ") if '.' in filename_or_title else filename_or_title,
        "authors": [],
        "year": "",
        "publisher": "",
        "keywords": [],
        "category": "Original Research",
        "type": "Literature"
    }

    # Ekstraksi Tahun
    year_match = re.search(r'\b(19|20)\d{2}\b', text_str[:5000])
    if year_match:
        metadata["year"] = year_match.group(0)

    # Heuristik Penerbit Sederhana
    publishers = ["Elsevier", "Springer", "IEEE", "MDPI", "Nature", "Science", "Wiley", "Taylor & Francis", "ACM", "Frontiers", "Sage", "Medium", "BBC", "CNN", "Wikipedia"]
    for pub in publishers:
        if pub.lower() in text_str[:10000].lower():
            metadata["publisher"] = pub
            break

    return metadata

def process_extracted_text(full_text, title, base_metadata=None):
    # Pastikan input adalah string
    text_str = str(full_text)
        
    # 1. Batasi total teks (200.000 karakter)
    limit_total = 200000
    limited_text = text_str[:limit_total]

    # 2. Gabungkan Heuristik atau Metadata dari API
    if base_metadata:
        metadata = {**extract_metadata_heuristics(limited_text, title), **base_metadata}
    else:
        metadata = extract_metadata_heuristics(limited_text, title)
    
    # 3. Snippet untuk AI (Groq/Gemini) - 7500 karakter
    ai_snippet = limited_text[:7500]
    
    # 4. Split teks ke dalam 10 chunks (masing-masing 20.000 karakter)
    chunk_size = 20000
    chunks = [limited_text[i:i+chunk_size] for i in range(0, len(limited_text), chunk_size)][:10]

    return {
        **metadata,
        "aiSnippet": ai_snippet,
        "chunks": chunks
    }

@app.route('/api/extract', methods=['POST'])
def extract():
    try:
        # HANDLING URL EXTRACTION (JSON Payload)
        if request.is_json:
            data = request.get_json()
            url = data.get('url')
            if not url:
                return jsonify({"status": "error", "message": "No URL provided"}), 400
            
            # --- TAHAP 1: CEK DOI (AKADEMIK) ---
            doi = extract_doi(url)
            academic_meta = None
            if doi:
                academic_meta = fetch_from_crossref(doi)
            
            # --- TAHAP 2: SCRAPING (GENERAL FALLBACK) ---
            try:
                session = requests.Session()
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/121.0.0.0 Safari/537.36',
                    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
                    'Accept-Language': 'en-US,en;q=0.9',
                    'Cache-Control': 'no-cache',
                    'Pragma': 'no-cache'
                }
                
                response = session.get(url, headers=headers, timeout=20, allow_redirects=True)
                response.raise_for_status()
                
                # Gunakan response.text (string) secara eksplisit
                html_str = str(response.text)
                
                # Alur: Readability
                doc = ReadabilityDoc(html_str)
                summary_html = doc.summary()
                title = academic_meta.get('title') if academic_meta else doc.short_title()
                
                # Alur: BeautifulSoup
                soup = BeautifulSoup(summary_html, 'lxml')
                for s in soup(['script', 'style', 'nav', 'header', 'footer', 'aside', 'iframe', 'button', 'input']):
                    s.decompose()
                
                full_text = clean_text(soup.get_text(separator=' '))
                
                if not full_text.strip():
                    if academic_meta:
                        # Jika scraping diblokir tapi Metadata API ada, tetap kirim metadata
                        return jsonify({"status": "success", "data": process_extracted_text("Content extraction blocked by provider, but metadata retrieved via DOI.", title, academic_meta)})
                    return jsonify({"status": "error", "message": "Content extraction failed. The page might be protected."}), 422
                
                result_data = process_extracted_text(full_text, title, academic_meta)
                return jsonify({"status": "success", "data": result_data})

            except requests.exceptions.HTTPError as e:
                # Jika scraping gagal tapi kita punya Metadata DOI, gunakan itu
                if academic_meta:
                    return jsonify({"status": "success", "data": process_extracted_text("Full text extraction failed due to access restrictions, but metadata was found.", academic_meta.get('title'), academic_meta)})
                
                status_code = e.response.status_code if e.response else 500
                return jsonify({"status": "error", "message": f"Website blocked access (Error {status_code})."}), status_code
            except Exception as e:
                return jsonify({"status": "error", "message": f"Web Extraction Error: {str(e)}"}), 500

        # HANDLING FILE EXTRACTION (Multipart Form Data)
        if 'file' not in request.files:
            return jsonify({"status": "error", "message": "No file part"}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({"status": "error", "message": "No file selected"}), 400

        filename_lower = file.filename.lower()
        file_bytes = file.read()
        f = io.BytesIO(file_bytes)
        full_text = ""

        if filename_lower.endswith('.pdf'):
            reader = PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text: full_text += page_text + "\n"
        elif filename_lower.endswith('.docx'):
            doc = Document(f)
            full_text = "\n".join([para.text for para in doc.paragraphs])
        elif filename_lower.endswith('.pptx'):
            prs = Presentation(f)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text: full_text += shape.text + "\n"
        elif filename_lower.endswith('.xlsx'):
            wb = openpyxl.load_workbook(f, data_only=True)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    full_text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        elif filename_lower.endswith(('.txt', '.md', '.csv')):
            full_text = file_bytes.decode('utf-8', errors='ignore')
        else:
            return jsonify({"status": "error", "message": "Unsupported file format."}), 400

        if not full_text.strip():
            return jsonify({"status": "error", "message": "Document is empty."}), 422

        result_data = process_extracted_text(clean_text(full_text), file.filename)
        return jsonify({"status": "success", "data": result_data})
        
    except Exception as e:
        return jsonify({"status": "error", "message": f"Server Error: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(port=5000)
